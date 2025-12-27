#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
pptx_build.py - Markdown原稿からPPTX（スライド）生成
"""

import sys
import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from utils import ProjectPaths, logger


class PPTXBuilder:
    """PPTX生成クラス（デザイン改善版）"""
    
    # テーマカラー
    THEME_COLOR_PRIMARY = RGBColor(63, 81, 181)  # Indigo
    THEME_COLOR_ACCENT = RGBColor(33, 150, 243)  # Blue
    THEME_COLOR_SUCCESS = RGBColor(76, 175, 80)  # Green
    THEME_COLOR_WARNING = RGBColor(255, 152, 0)  # Orange
    THEME_COLOR_DANGER = RGBColor(244, 67, 54)  # Red
    
    def __init__(self, output_dir: Path):
        self.output_dir = output_dir
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.figures_dir = self.output_dir.parent / "assets" / "figs"
    
    def apply_theme_to_shape(self, shape, is_title=False):
        """テーマカラーを適用"""
        if hasattr(shape, 'text_frame'):
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if is_title:
                        run.font.size = Pt(32)
                        run.font.bold = True
                        run.font.color.rgb = self.THEME_COLOR_PRIMARY
                    else:
                        run.font.size = Pt(18)
                        run.font.color.rgb = RGBColor(50, 50, 50)
    
    def create_title_slide(self, prs: Presentation, title: str, subtitle: str = ""):
        """タイトルスライド作成（改善版）"""
        # ブランクレイアウトを使用
        blank_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_layout)
        
        # 背景グラデーション風の矩形
        shapes = slide.shapes
        bg_rect = shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(0),
            prs.slide_width, prs.slide_height
        )
        bg_rect.fill.solid()
        bg_rect.fill.fore_color.rgb = RGBColor(240, 245, 250)
        bg_rect.line.fill.background()
        
        # タイトル
        title_box = shapes.add_textbox(
            Inches(1), Inches(2.5),
            Inches(8), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        p = title_frame.paragraphs[0]
        p.text = title
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = self.THEME_COLOR_PRIMARY
        
        # サブタイトル
        if subtitle:
            subtitle_box = shapes.add_textbox(
                Inches(1), Inches(4.2),
                Inches(8), Inches(0.8)
            )
            subtitle_frame = subtitle_box.text_frame
            p = subtitle_frame.paragraphs[0]
            p.text = subtitle
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(24)
            p.font.color.rgb = self.THEME_COLOR_ACCENT
        
        return slide
    
    def create_content_slide(self, prs: Presentation, title: str, content: str, chapter_num: int = 0):
        """コンテンツスライド作成（改善版）"""
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        shapes = slide.shapes
        
        # 背景
        bg_rect = shapes.add_shape(
            1, Inches(0), Inches(0),
            prs.slide_width, prs.slide_height
        )
        bg_rect.fill.solid()
        bg_rect.fill.fore_color.rgb = RGBColor(255, 255, 255)
        bg_rect.line.fill.background()
        
        # ヘッダー帯
        header_rect = shapes.add_shape(
            1, Inches(0), Inches(0),
            prs.slide_width, Inches(0.8)
        )
        header_rect.fill.solid()
        header_rect.fill.fore_color.rgb = self.THEME_COLOR_PRIMARY
        header_rect.line.fill.background()
        
        # タイトル
        title_box = shapes.add_textbox(
            Inches(0.5), Inches(0.15),
            Inches(8.5), Inches(0.5)
        )
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        # コンテンツエリア
        content_box = shapes.add_textbox(
            Inches(0.8), Inches(1.2),
            Inches(8.4), Inches(5.8)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        # コンテンツを整形
        lines = content.split('\n')
        line_count = 0
        for line in lines:
            if line_count >= 15:  # 最大15行
                break
            
            line = line.strip()
            if not line:
                continue
            
            # 箇条書き記号の検出
            if line.startswith('- ') or line.startswith('* '):
                line = line[2:]
                is_bullet = True
            elif line.startswith('✓ ') or line.startswith('✅'):
                line = line[2:]
                is_bullet = True
                color = self.THEME_COLOR_SUCCESS
            elif line.startswith('❌') or line.startswith('✗'):
                line = line[2:]
                is_bullet = True
                color = self.THEME_COLOR_DANGER
            else:
                is_bullet = False
                color = RGBColor(50, 50, 50)
            
            if line_count == 0:
                p = content_frame.paragraphs[0]
            else:
                p = content_frame.add_paragraph()
            
            p.text = line[:120]  # 最大120文字
            p.font.size = Pt(16)
            p.font.color.rgb = color if 'color' in locals() else RGBColor(50, 50, 50)
            
            if is_bullet:
                p.level = 1
            
            line_count += 1
        
        # フッター
        footer_box = shapes.add_textbox(
            Inches(8), Inches(7.0),
            Inches(1.5), Inches(0.3)
        )
        footer_frame = footer_box.text_frame
        p = footer_frame.paragraphs[0]
        p.text = f"第{chapter_num}章"
        p.alignment = PP_ALIGN.RIGHT
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(150, 150, 150)
        
        return slide
    
    def create_figure_slide(self, prs: Presentation, title: str, figure_name: str, chapter_num: int = 0):
        """図表スライド作成"""
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        shapes = slide.shapes
        
        # 背景
        bg_rect = shapes.add_shape(
            1, Inches(0), Inches(0),
            prs.slide_width, prs.slide_height
        )
        bg_rect.fill.solid()
        bg_rect.fill.fore_color.rgb = RGBColor(250, 250, 250)
        bg_rect.line.fill.background()
        
        # タイトル
        title_box = shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(9), Inches(0.6)
        )
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = title
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.THEME_COLOR_PRIMARY
        
        # 図の挿入
        figure_path = self.figures_dir / f"{figure_name}.png"
        if figure_path.exists():
            try:
                shapes.add_picture(
                    str(figure_path),
                    Inches(1.5), Inches(1.5),
                    width=Inches(7), height=Inches(5)
                )
            except Exception as e:
                logger.warning(f"図の挿入に失敗: {figure_name} - {e}")
        
        return slide
    
    def build_pptx(self, volume_name: str, manuscript_dir: Path):
        """PPTXを生成"""
        output_file = self.output_dir / f"{volume_name}.pptx"
        
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # 表紙
        if "vol1" in volume_name:
            title = "BIM利用技術者試験\n2級対応 教科書"
            subtitle = "Building Information Modeling for 2級 Certification"
        else:
            title = "BIM利用技術者試験\n準1級対応 教科書"
            subtitle = "Building Information Modeling for 準1級 Certification"
        
        self.create_title_slide(prs, title, subtitle)
        
        # 各章からスライド作成
        chapter_files = sorted(manuscript_dir.glob("chapter_*.md"))
        for idx, chapter_file in enumerate(chapter_files, start=1):
            logger.info(f"  処理中: {chapter_file.name}")
            content = chapter_file.read_text(encoding='utf-8')
            
            # 章タイトル抽出
            title_match = re.search(r'^# (.+)$', content, re.MULTILINE)
            if title_match:
                chapter_title = title_match.group(1)
            else:
                chapter_title = chapter_file.stem
            
            # 章タイトルスライド
            self.create_title_slide(prs, f"第{idx}章", chapter_title)
            
            # セクションごとにスライド作成
            sections = re.split(r'####\s+', content)[1:8]  # 最大8セクション
            for section in sections:
                lines = section.split('\n')
                section_title = lines[0].strip() if lines else "内容"
                
                # コンテンツ抽出（図表参照を除外）
                section_lines = []
                for line in lines[1:]:
                    if line.strip().startswith('![FIG:'):
                        # 図表参照を検出
                        fig_match = re.search(r'!\[FIG:(\w+)\]', line)
                        if fig_match:
                            fig_name = fig_match.group(1)
                            # 図表スライドを作成
                            self.create_figure_slide(
                                prs,
                                f"{chapter_title} - 図表",
                                fig_name,
                                idx
                            )
                    elif not line.strip().startswith('```'):
                        section_lines.append(line)
                
                section_content = '\n'.join(section_lines[:25])
                
                if section_content.strip():
                    self.create_content_slide(
                        prs,
                        f"{chapter_title} - {section_title}",
                        section_content,
                        idx
                    )
        
        # 保存
        prs.save(str(output_file))
        logger.info(f"✅ PPTX生成完了: {output_file}")


def build_all_pptx():
    """PPTXを生成（build.pyから呼ばれる）"""
    paths = ProjectPaths()
    builder = PPTXBuilder(paths.dist)
    
    logger.info("=" * 70)
    logger.info("PPTX生成開始")
    logger.info("=" * 70)
    
    # VOL1
    logger.info("VOL1 (2級対応) PPTX生成中...")
    builder.build_pptx("vol1_2kyu", paths.vol1_dir)
    
    # VOL2
    logger.info("VOL2 (準1級対応) PPTX生成中...")
    builder.build_pptx("vol2_jun1kyu", paths.vol2_dir)
    
    logger.info("=" * 70)
    logger.info("✨ PPTX生成完了")
    logger.info(f"📊 VOL1: {paths.dist / 'vol1_2kyu.pptx'}")
    logger.info(f"📊 VOL2: {paths.dist / 'vol2_jun1kyu.pptx'}")
    logger.info("=" * 70)


if __name__ == "__main__":
    try:
        build_all_pptx()
        logger.info("✅ PPTX生成完了")
    except Exception as e:
        logger.error(f"❌ エラーが発生しました: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)
